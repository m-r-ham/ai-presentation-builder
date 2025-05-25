#!/usr/bin/env python3
"""
Extract real slide images from PowerPoint files for AI training
Uses LibreOffice to convert PPTX to PDF, then PDF to images
"""

import json
import sys
import os
import subprocess
import tempfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pdf2image import convert_from_path
    from PIL import Image
except ImportError:
    print("Please install required packages:")
    print("pip install python-pptx pdf2image pillow")
    print("Also install LibreOffice and poppler-utils")
    sys.exit(1)

def extract_text_from_shape(shape):
    """Extract text from a shape if it has text."""
    if hasattr(shape, "text"):
        return shape.text.strip()
    return ""

def extract_slide_metadata(slide, slide_number):
    """Extract metadata from a single slide."""
    slide_data = {
        "id": f"slide-{slide_number:03d}",
        "slide_number": slide_number,
        "title": "",
        "content": "",
        "bullet_points": [],
        "has_image": False,
        "has_chart": False,
        "has_table": False,
        "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown",
        "text_content": [],
        "shape_count": len(slide.shapes),
        "design_elements": {
            "text_blocks": 0,
            "visual_elements": 0,
            "color_usage": "unknown",
            "layout_complexity": "unknown"
        }
    }
    
    # Extract text from all shapes
    text_blocks = 0
    visual_elements = 0
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = extract_text_from_shape(shape)
            if text:
                slide_data["text_content"].append(text)
                text_blocks += 1
                
                # Try to identify title (usually first or short text)
                if not slide_data["title"] and len(text) < 100:
                    slide_data["title"] = text
                elif not slide_data["content"]:
                    slide_data["content"] = text
                else:
                    slide_data["bullet_points"].append(text)
        
        # Count visual elements
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            slide_data["has_image"] = True
            visual_elements += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            slide_data["has_chart"] = True
            visual_elements += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            slide_data["has_table"] = True
            visual_elements += 1
    
    # Set design elements
    slide_data["design_elements"]["text_blocks"] = text_blocks
    slide_data["design_elements"]["visual_elements"] = visual_elements
    
    # Determine layout complexity
    total_elements = text_blocks + visual_elements
    if total_elements <= 3:
        slide_data["design_elements"]["layout_complexity"] = "simple"
    elif total_elements <= 6:
        slide_data["design_elements"]["layout_complexity"] = "moderate"
    else:
        slide_data["design_elements"]["layout_complexity"] = "complex"
    
    # Fallback title
    if not slide_data["title"] and slide_data["text_content"]:
        slide_data["title"] = slide_data["text_content"][0][:50] + "..." if len(slide_data["text_content"][0]) > 50 else slide_data["text_content"][0]
    
    # Combine remaining text as content
    if not slide_data["content"] and len(slide_data["text_content"]) > 1:
        slide_data["content"] = " ".join(slide_data["text_content"][1:])
    
    return slide_data

def convert_pptx_to_images(pptx_path, output_dir):
    """Convert PowerPoint to images using LibreOffice -> PDF -> Images"""
    try:
        # Create temporary directory for PDF
        with tempfile.TemporaryDirectory() as temp_dir:
            print(f"Converting {pptx_path} to PDF...")
            
            # Convert PPTX to PDF using LibreOffice
            pdf_path = os.path.join(temp_dir, "presentation.pdf")
            
            # Try LibreOffice command
            try:
                subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', temp_dir, pptx_path
                ], check=True, capture_output=True)
                
                # Find the generated PDF file
                pdf_files = list(Path(temp_dir).glob("*.pdf"))
                if pdf_files:
                    pdf_path = str(pdf_files[0])
                else:
                    raise FileNotFoundError("PDF not generated")
                    
            except (subprocess.CalledProcessError, FileNotFoundError) as e:
                print(f"LibreOffice conversion failed: {e}")
                print("Trying alternative method...")
                
                # Try soffice command as alternative
                subprocess.run([
                    'soffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', temp_dir, pptx_path
                ], check=True, capture_output=True)
                
                pdf_files = list(Path(temp_dir).glob("*.pdf"))
                if pdf_files:
                    pdf_path = str(pdf_files[0])
                else:
                    raise FileNotFoundError("PDF not generated by soffice either")
            
            print(f"PDF created: {pdf_path}")
            
            # Convert PDF pages to images
            print("Converting PDF pages to images...")
            images = convert_from_path(pdf_path, dpi=150, fmt='PNG')
            
            # Save images
            image_paths = []
            for i, image in enumerate(images, 1):
                image_filename = f"slide_{i:03d}.png"
                image_path = os.path.join(output_dir, image_filename)
                
                # Resize image for web display (max 800px width)
                max_width = 800
                if image.width > max_width:
                    ratio = max_width / image.width
                    new_height = int(image.height * ratio)
                    image = image.resize((max_width, new_height), Image.Resampling.LANCZOS)
                
                image.save(image_path, 'PNG', optimize=True)
                image_paths.append(image_filename)
                print(f"Saved: {image_filename}")
            
            return image_paths
            
    except Exception as e:
        print(f"Error converting PPTX to images: {e}")
        return []

def extract_presentation_data(pptx_path, output_dir):
    """Extract presentation data and create images + JSON files"""
    try:
        # Create output directories
        images_dir = os.path.join(output_dir, "images")
        metadata_dir = os.path.join(output_dir, "metadata")
        os.makedirs(images_dir, exist_ok=True)
        os.makedirs(metadata_dir, exist_ok=True)
        
        print(f"Processing presentation: {pptx_path}")
        
        # Convert slides to images
        image_filenames = convert_pptx_to_images(pptx_path, images_dir)
        
        if not image_filenames:
            print("Failed to extract slide images")
            return False
        
        # Extract metadata from PPTX
        presentation = Presentation(pptx_path)
        slides_metadata = []
        
        print(f"Extracting metadata for {len(presentation.slides)} slides...")
        
        for i, slide in enumerate(presentation.slides, 1):
            if i <= len(image_filenames):  # Only process slides we have images for
                slide_metadata = extract_slide_metadata(slide, i)
                slide_metadata["image_filename"] = image_filenames[i-1]
                slide_metadata["image_path"] = f"images/{image_filenames[i-1]}"
                
                # Save individual JSON file for each slide
                json_filename = f"slide_{i:03d}.json"
                json_path = os.path.join(metadata_dir, json_filename)
                
                with open(json_path, 'w', encoding='utf-8') as f:
                    json.dump(slide_metadata, f, indent=2, ensure_ascii=False)
                
                slides_metadata.append(slide_metadata)
                print(f"Created: {json_filename}")
        
        # Create master metadata file
        master_metadata = {
            "source_file": os.path.basename(pptx_path),
            "total_slides": len(slides_metadata),
            "extraction_date": str(Path(pptx_path).stat().st_mtime),
            "extraction_method": "libreoffice_pdf2image",
            "images_directory": "images",
            "metadata_directory": "metadata",
            "slides": slides_metadata
        }
        
        master_path = os.path.join(output_dir, "presentation_data.json")
        with open(master_path, 'w', encoding='utf-8') as f:
            json.dump(master_metadata, f, indent=2, ensure_ascii=False)
        
        print(f"\nExtraction complete!")
        print(f"- {len(image_filenames)} slide images saved to: {images_dir}")
        print(f"- {len(slides_metadata)} metadata files saved to: {metadata_dir}")
        print(f"- Master data file: {master_path}")
        
        return True
        
    except Exception as e:
        print(f"Error extracting presentation data: {e}")
        return False

def main():
    if len(sys.argv) != 3:
        print("Usage: python extract_pptx_images.py <input.pptx> <output_directory>")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_dir = sys.argv[2]
    
    if not os.path.exists(input_file):
        print(f"Error: Input file {input_file} not found")
        sys.exit(1)
    
    success = extract_presentation_data(input_file, output_dir)
    if success:
        print("✅ Slide extraction completed successfully!")
    else:
        print("❌ Slide extraction failed!")
        sys.exit(1)

if __name__ == "__main__":
    main()