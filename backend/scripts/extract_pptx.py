#!/usr/bin/env python3
"""
Extract slide content and images from PowerPoint files for AI training
"""

import json
import sys
import os
from pathlib import Path
import io
from PIL import Image, ImageDraw, ImageFont

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches
except ImportError:
    print("Please install python-pptx and pillow: pip install python-pptx pillow")
    sys.exit(1)

def extract_text_from_shape(shape):
    """Extract text from a shape if it has text."""
    if hasattr(shape, "text"):
        return shape.text.strip()
    return ""

def create_slide_thumbnail(slide, slide_number, output_dir):
    """Create a visual representation of the slide."""
    try:
        # Create a simple visual representation based on slide content
        img_width, img_height = 800, 600
        img = Image.new('RGB', (img_width, img_height), color='white')
        draw = ImageDraw.Draw(img)
        
        # Try to load a system font, fallback to default
        try:
            title_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 24)
            content_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 16)
            small_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 12)
        except:
            title_font = ImageFont.load_default()
            content_font = ImageFont.load_default()
            small_font = ImageFont.load_default()
        
        y_offset = 50
        
        # Draw slide number
        draw.text((20, 20), f"Slide {slide_number}", fill='gray', font=small_font)
        
        # Process shapes and create visual representation
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                text = shape.text.strip()
                
                # Determine if this is likely a title (first large text or short text)
                if y_offset == 50 and len(text) < 100:
                    # Draw as title
                    wrapped_text = wrap_text(text, 35)
                    for line in wrapped_text[:2]:  # Max 2 lines for title
                        draw.text((50, y_offset), line, fill='black', font=title_font)
                        y_offset += 35
                    y_offset += 20
                else:
                    # Draw as content
                    wrapped_text = wrap_text(text, 60)
                    for line in wrapped_text[:8]:  # Max 8 lines for content
                        draw.text((50, y_offset), f"• {line}", fill='darkblue', font=content_font)
                        y_offset += 25
                    y_offset += 15
            
            # Draw visual indicators for charts, images, tables
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                draw.rectangle([50, y_offset, 200, y_offset + 80], outline='blue', width=2)
                draw.text((60, y_offset + 30), "[CHART]", fill='blue', font=content_font)
                y_offset += 100
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                draw.rectangle([50, y_offset, 200, y_offset + 80], outline='green', width=2)
                draw.text((60, y_offset + 30), "[IMAGE]", fill='green', font=content_font)
                y_offset += 100
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                draw.rectangle([50, y_offset, 200, y_offset + 80], outline='purple', width=2)
                draw.text((60, y_offset + 30), "[TABLE]", fill='purple', font=content_font)
                y_offset += 100
        
        # Add layout name at bottom
        layout_name = slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown"
        draw.text((20, img_height - 30), f"Layout: {layout_name}", fill='gray', font=small_font)
        
        # Save the image
        image_filename = f"slide_{slide_number:03d}.png"
        image_path = os.path.join(output_dir, image_filename)
        img.save(image_path, 'PNG')
        
        return image_filename
        
    except Exception as e:
        print(f"Error creating thumbnail for slide {slide_number}: {e}")
        return None

def wrap_text(text, width):
    """Wrap text to specified width."""
    words = text.split()
    lines = []
    current_line = []
    
    for word in words:
        if len(' '.join(current_line + [word])) <= width:
            current_line.append(word)
        else:
            if current_line:
                lines.append(' '.join(current_line))
                current_line = [word]
            else:
                lines.append(word)
    
    if current_line:
        lines.append(' '.join(current_line))
    
    return lines

def extract_slide_content(slide, slide_number, output_dir):
    """Extract content from a single slide and create thumbnail."""
    slide_data = {
        "id": f"template-{slide_number}",
        "slide_number": slide_number,
        "title": "",
        "content": "",
        "bullet_points": [],
        "has_image": False,
        "has_chart": False,
        "has_table": False,
        "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown",
        "text_content": [],
        "shape_count": len(slide.shapes),
        "image_filename": None
    }
    
    # Extract text from all shapes
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = extract_text_from_shape(shape)
            if text:
                slide_data["text_content"].append(text)
                
                # Try to identify title (usually first or largest text)
                if not slide_data["title"] and len(text) < 100:
                    slide_data["title"] = text
                elif not slide_data["content"]:
                    slide_data["content"] = text
                else:
                    slide_data["bullet_points"].append(text)
        
        # Check for images
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            slide_data["has_image"] = True
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            slide_data["has_chart"] = True
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            slide_data["has_table"] = True
    
    # Fallback: if no title found, use first text content
    if not slide_data["title"] and slide_data["text_content"]:
        slide_data["title"] = slide_data["text_content"][0][:50] + "..." if len(slide_data["text_content"][0]) > 50 else slide_data["text_content"][0]
    
    # Combine all text as content if no specific content found
    if not slide_data["content"] and len(slide_data["text_content"]) > 1:
        slide_data["content"] = " ".join(slide_data["text_content"][1:])
    
    # Create slide thumbnail
    image_filename = create_slide_thumbnail(slide, slide_number, output_dir)
    slide_data["image_filename"] = image_filename
    
    return slide_data

def extract_powerpoint_slides(pptx_path, output_path):
    """Extract all slides from PowerPoint file."""
    try:
        presentation = Presentation(pptx_path)
        slides_data = []
        
        # Create images directory
        images_dir = os.path.dirname(output_path) + "/images"
        os.makedirs(images_dir, exist_ok=True)
        
        print(f"Processing {len(presentation.slides)} slides from {pptx_path}")
        print(f"Images will be saved to: {images_dir}")
        
        for i, slide in enumerate(presentation.slides, 1):
            slide_data = extract_slide_content(slide, i, images_dir)
            
            # Only include slides with actual content
            if slide_data["title"] or slide_data["content"] or slide_data["text_content"]:
                slides_data.append(slide_data)
                print(f"Extracted slide {i}: {slide_data['title'][:30]}...")
        
        # Save to JSON file
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump({
                "source_file": os.path.basename(pptx_path),
                "total_slides": len(slides_data),
                "extraction_date": str(Path(pptx_path).stat().st_mtime),
                "images_directory": "images",
                "slides": slides_data
            }, f, indent=2, ensure_ascii=False)
        
        print(f"Extracted {len(slides_data)} slides to {output_path}")
        print(f"Created {len(slides_data)} slide images in {images_dir}")
        return slides_data
        
    except Exception as e:
        print(f"Error processing {pptx_path}: {e}")
        return []

def main():
    if len(sys.argv) != 3:
        print("Usage: python extract_pptx.py <input.pptx> <output.json>")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    
    if not os.path.exists(input_file):
        print(f"Error: Input file {input_file} not found")
        sys.exit(1)
    
    extract_powerpoint_slides(input_file, output_file)

if __name__ == "__main__":
    main()